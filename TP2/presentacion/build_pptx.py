#!/usr/bin/env python3
"""Arma el .pptx que se usa PARA PRESENTAR, con las animaciones embebidas.

La catedra pidio dos archivos distintos: el PDF que se entrega lleva un
fotograma fijo y el link a la animacion, y el que se presenta tiene que
reproducir la animacion sin salir de la diapositiva. PowerPoint reproduce GIF
animados en modo presentacion, asi que la version en vivo es un .pptx.

El .pptx NO se re-maqueta a mano: se compila `presentacion.tex` en modo `vivo`
(donde cada animacion deja un rectangulo marcador en vez del fotograma), se
rasteriza cada pagina y se pega esa imagen a sangre en una diapositiva. Asi el
diseno es identico al del PDF, incluida la barra de secciones de `miniframes`.
Encima de cada rectangulo marcador -- localizado por color en la pagina
rasterizada, no por coordenadas escritas a mano -- se coloca el GIF.

    py presentacion/build_pptx.py            # desde TP2/
    py presentacion/build_pptx.py --dpi 220  # mas resolucion, .pptx mas pesado
"""

import argparse
import subprocess
import sys
from pathlib import Path

import fitz  # PyMuPDF
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Inches

TP2_DIR = Path(__file__).resolve().parent.parent
PRES_DIR = TP2_DIR / "presentacion"
TEX_NAME = "presentacion.tex"
LIVE_JOBNAME = "presentacion_vivo"
ANIM_DIR = TP2_DIR / "data" / "plots" / "animation"

# Mismo valor que \definecolor{huecoanim} en presentacion.tex.
MARKER_RGB = (255, 0, 255)
MARKER_TOL = 12

# Diapositiva 16:9, la misma relacion de aspecto que `aspectratio=169`.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Orden en que se consumen los GIF: de arriba a abajo y de izquierda a derecha
# dentro de cada diapositiva, que es el orden en que aparecen los \animacion.
GIF_SEQUENCE = [
    "animation_vicsek_rho2_ordenado.gif",
    "animation_vicsek_rho2_desordenado.gif",
    "animation_voter_rho2_ordenado.gif",
    "animation_voter_rho2_desordenado.gif",
]


def compile_live_pdf(passes: int = 2) -> Path:
    """Compila la variante `vivo` de la presentacion.

    Dos pasadas porque `miniframes` necesita el .aux de la corrida anterior
    para dibujar la barra de secciones con la cantidad correcta de puntos.
    """
    for i in range(passes):
        proc = subprocess.run(
            ["pdflatex", "-interaction=nonstopmode", f"-jobname={LIVE_JOBNAME}",
             rf"\def\modo{{vivo}}\input{{{TEX_NAME}}}"],
            cwd=PRES_DIR, capture_output=True, text=True,
        )
        if proc.returncode != 0:
            tail = "\n".join(proc.stdout.splitlines()[-25:])
            raise RuntimeError(f"pdflatex fallo en la pasada {i + 1}:\n{tail}")
    return PRES_DIR / f"{LIVE_JOBNAME}.pdf"


def marker_mask(page_rgb: np.ndarray) -> np.ndarray:
    """Mascara booleana de los pixeles del color marcador."""
    return np.all(np.abs(page_rgb.astype(int) - MARKER_RGB) <= MARKER_TOL, axis=2)


def marker_boxes(mask: np.ndarray) -> list[tuple[int, int, int, int]]:
    """Rectangulos marcadores de una pagina rasterizada, en pixeles.

    Devuelve una lista de (x0, y0, x1, y1) ordenada de izquierda a derecha. Se
    separan por columnas de pixeles marcados: los huecos de una diapositiva
    estan siempre uno al lado del otro, nunca apilados.
    """
    if not mask.any():
        return []

    columns = np.nonzero(mask.any(axis=0))[0]
    boxes = []
    start = prev = columns[0]
    for col in columns[1:]:
        if col > prev + 1:                     # se corto la mancha
            boxes.append((start, prev))
            start = col
        prev = col
    boxes.append((start, prev))

    out = []
    for x0, x1 in boxes:
        rows = np.nonzero(mask[:, x0:x1 + 1].any(axis=1))[0]
        out.append((x0, int(rows[0]), x1, int(rows[-1])))
    return out


def fit_box(box_px, page_px, gif_path: Path):
    """Ubicacion del GIF dentro del hueco, en EMU, conservando su aspecto.

    El GIF y el fotograma no tienen exactamente la misma relacion de aspecto
    (los genera matplotlib con figuras distintas), asi que el GIF se centra
    dentro del hueco en vez de estirarse.
    """
    x0, y0, x1, y1 = box_px
    page_w, page_h = page_px
    box_w = (x1 - x0 + 1) / page_w * SLIDE_W
    box_h = (y1 - y0 + 1) / page_h * SLIDE_H
    left = x0 / page_w * SLIDE_W
    top = y0 / page_h * SLIDE_H

    with Image.open(gif_path) as img:
        aspect = img.width / img.height

    if box_w / box_h > aspect:                  # sobra ancho
        height, width = box_h, box_h * aspect
    else:
        width, height = box_w, box_w / aspect
    left += (box_w - width) / 2
    top += (box_h - height) / 2
    return Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height))


def build(pdf_path: Path, out_path: Path, dpi: int, page_dir: Path) -> int:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]                # layout sin placeholders

    page_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(pdf_path)
    pending = list(GIF_SEQUENCE)
    placed = 0

    for index, page in enumerate(doc):
        pix = page.get_pixmap(dpi=dpi)
        rgb = (np.frombuffer(pix.samples, dtype=np.uint8)
                 .reshape(pix.height, pix.width, pix.n)[:, :, :3]).copy()
        mask = marker_mask(rgb)
        # El marcador se borra del fondo antes de guardar la pagina. Se blanquea
        # el rectangulo entero con un margen, no solo los pixeles que matchean:
        # el borde del rectangulo sale antialiaseado contra el fondo y esos
        # pixeles intermedios sobrevivirian al filtro de color, dejando un
        # contorno de colores alrededor de la animacion.
        boxes = marker_boxes(mask)
        margin = max(2, round(dpi / 50))
        for x0, y0, x1, y1 in boxes:
            rgb[max(0, y0 - margin):y1 + 1 + margin,
                max(0, x0 - margin):x1 + 1 + margin] = 255

        png = page_dir / f"pagina{index + 1:02d}.png"
        Image.fromarray(rgb).save(png)

        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, width=SLIDE_W, height=SLIDE_H)

        for box in boxes:
            if not pending:
                raise RuntimeError(
                    f"la pagina {index + 1} tiene un hueco de animacion de mas: "
                    f"GIF_SEQUENCE lista {len(GIF_SEQUENCE)} archivos"
                )
            gif = ANIM_DIR / pending.pop(0)
            if not gif.exists():
                raise FileNotFoundError(f"falta {gif}: correr `make animation` primero")
            left, top, width, height = fit_box(box, (pix.width, pix.height), gif)
            slide.shapes.add_picture(str(gif), left, top, width=width, height=height)
            placed += 1

    if pending:
        raise RuntimeError(f"quedaron GIF sin ubicar: {pending}")

    try:
        prs.save(out_path)
    except PermissionError:
        # Windows bloquea el archivo mientras PowerPoint lo tiene abierto, que
        # es justo lo que pasa cuando uno revisa el .pptx y lo regenera.
        raise SystemExit(
            f"error: {out_path} esta abierto en PowerPoint (o bloqueado por otro "
            f"proceso). Cerralo y volve a correr."
        )
    return placed


def main():
    parser = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    parser.add_argument("--dpi", type=int, default=200,
                        help="resolucion del rasterizado de cada pagina (default 200)")
    parser.add_argument("--out", type=Path, default=PRES_DIR / "presentacion_vivo.pptx")
    parser.add_argument("--skip-latex", action="store_true",
                        help="reusar presentacion_vivo.pdf en vez de recompilarlo")
    args = parser.parse_args()

    pdf = PRES_DIR / f"{LIVE_JOBNAME}.pdf"
    if not args.skip_latex:
        pdf = compile_live_pdf()
    elif not pdf.exists():
        sys.exit(f"error: no existe {pdf} y se pidio --skip-latex")

    page_dir = PRES_DIR / "_paginas_vivo"
    placed = build(pdf, args.out, args.dpi, page_dir)
    size_mb = args.out.stat().st_size / 1e6
    print(f"{args.out}  ({placed} animaciones embebidas, {size_mb:.1f} MB)")
    print("Abrirlo con PowerPoint y presentar con F5: los GIF se reproducen solos.")


if __name__ == "__main__":
    main()
